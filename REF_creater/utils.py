from docx import Document
import PyPDF2
from pptx import Presentation

def load_report_text(uploaded_file):
    filename = uploaded_file.name.lower()

    if filename.endswith(".txt"):
        return uploaded_file.read().decode("utf-8")

    elif filename.endswith(".docx"):
        doc = Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(uploaded_file)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return text.strip()

    elif filename.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    else:
        return ""
